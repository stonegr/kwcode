"""
OfficeHandler expert: generates Office documents (docx/xlsx/pptx) via
LLM-generated Python scripts executed with run_bash.

Pipeline: detect type → select scene prompt → LLM generates script → execute → verify file exists.
"""

import logging
import os
import re
import tempfile
import time

from kaiwu.core.context import TaskContext

logger = logging.getLogger(__name__)

# ── Scene system prompts (distilled from cl-v2 scenes) ──────────────────

XLSX_SYSTEM = """\
你是Excel文档生成专家。用openpyxl生成专业的.xlsx文件。

## 强制规则
1. 用Excel公式（=SUM等），不要Python里算好再硬编码值
2. 必须应用下方样式，禁止生成无格式白底表格

## openpyxl样式规范（必须使用）
配色方案（商务默认）：
  表头背景 #1B2A4A，表头文字白色粗体居中
  交替行色 #F0F4FA（斑马纹）
  边框浅灰细线 #D1D5DB
  汇总行粗体 + 双线上边框 + 背景 #E8EBF2

必须设置：
  ws.freeze_panes = 'A2'（冻结表头）
  列宽自适应（column_dimensions.width >= 15）
  汇总行用Excel公式 =SUM(...)
  数字格式：货币用 '¥#,##0.00'，百分比用 '0.0%'，日期用 'YYYY-MM-DD'
  行高：表头行高28

## 脚本要求
- 开头加 Windows UTF-8 编码头（见下方模板）
- 只输出完整可执行的Python脚本
- 不要解释，不要markdown代码块标记
- 脚本末尾 print 输出文件路径"""

PPTX_SYSTEM = """\
你是PPT生成专家。用python-pptx生成专业演示文稿。

## 强制规则
1. 用 slide_layouts[6]（空白版式）+ 手动添加元素，不用默认占位符
2. 必须应用配色方案，禁止白底黑字默认样式
3. 每页必须有视觉元素（色块/图标emoji/表格）
4. 深-浅-深三明治结构：标题页/结尾页深色背景，内容页浅色背景
5. 绝对禁止标题下加装饰横线

## 只允许使用以下 import（禁止编造不存在的模块）
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

## 关键API用法（必须严格遵守，不要猜测）
- 添加文本框：shape = slide.shapes.add_textbox(left, top, width, height)
- 设置文字：tf = shape.text_frame; tf.text = "内容"
- 设置字体：p = tf.paragraphs[0]; run = p.runs[0]; run.font.size = Pt(40)
  注意：font 在 run 上，不在 text_frame 上！
- 设置颜色：run.font.color.rgb = RGBColor(0x1B, 0x2A, 0x4A)
- 设置粗体：run.font.bold = True
- 添加新段落：p = tf.add_paragraph(); p.text = "新行"
- 设置对齐：p.alignment = PP_ALIGN.CENTER
- 设置背景色：
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1B, 0x2A, 0x4A)
- word_wrap：tf.word_wrap = True
- 添加色块矩形：slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
- 禁止：text_frame.font（不存在）、slide.background = xxx（不能直接赋值）
  禁止：prs.slides.add_slide()不传参数（必须传slide_layout）
  禁止：MSO_SHAPE.RECTANGULAR_ARROW等不存在的形状，只用RECTANGLE/ROUNDED_RECTANGLE/OVAL

## 配色方案（商务/汇报）
PRIMARY = RGBColor(0x1B, 0x2A, 0x4A)    # 深蓝主色
SECONDARY = RGBColor(0x2D, 0x5F, 0x8A)  # 中蓝辅色
ACCENT = RGBColor(0xE8, 0xA8, 0x38)     # 琥珀强调
BG_DARK = RGBColor(0x1B, 0x2A, 0x4A)    # 深色背景
BG_LIGHT = RGBColor(0xF5, 0xF7, 0xFA)   # 浅色背景
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)

## 字体规范
标题 36-44pt bold，正文 14-16pt，副标题 18-20pt
text_frame.word_wrap = True 必须设置

## 幻灯片尺寸
prs.slide_width = Cm(33.867)   # 16:9
prs.slide_height = Cm(19.05)

## 脚本要求
- 开头加 Windows UTF-8 编码头
- 只输出完整可执行的Python脚本
- 不要解释，不要markdown代码块标记
- 脚本末尾 print 输出文件路径"""

DOCX_SYSTEM = """\
你是Word文档生成专家。用python-docx生成专业文档。

## 强制规则
1. 必须用python-docx生成.docx文件
2. 严格按照下方模板代码的写法，只修改文字内容和结构
3. 不要发明API调用，只用模板中出现过的方法

## 完整模板（照这个格式写，只改内容）

```python
import sys
if sys.platform == "win32":
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

from docx import Document
from docx.shared import Pt, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH

doc = Document()

# 标题
h = doc.add_heading('文档标题', level=0)
h.alignment = WD_ALIGN_PARAGRAPH.CENTER

# 一级标题
doc.add_heading('一、第一部分', level=1)

# 正文段落（带首行缩进）
p = doc.add_paragraph('这是正文内容。')
p.paragraph_format.first_line_indent = Cm(0.74)

# 列表
doc.add_paragraph('要点一', style='List Bullet')
doc.add_paragraph('要点二', style='List Bullet')

# 表格
table = doc.add_table(rows=3, cols=3, style='Table Grid')
table.cell(0, 0).text = '列A'
table.cell(0, 1).text = '列B'
table.cell(0, 2).text = '列C'
table.cell(1, 0).text = '数据1'
table.cell(1, 1).text = '数据2'
table.cell(1, 2).text = '数据3'

# 落款（右对齐）
p = doc.add_paragraph()
p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
p.add_run('XX部门')
p = doc.add_paragraph()
p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
p.add_run('2026年4月28日')

doc.save('输出路径.docx')
print('输出路径.docx')
```

## 注意
- 只用 doc.add_heading / doc.add_paragraph / doc.add_table 这三个方法
- 首行缩进用 p.paragraph_format.first_line_indent = Cm(0.74)
- 列表用 style='List Bullet' 或 style='List Number'
- 不要用 run._element、qn、OxmlElement 等底层API
- 不要设置字体（默认字体即可）
- 只输出完整可执行的Python脚本，不要解释"""

# Windows UTF-8 header that must be prepended to generated scripts
_SCRIPT_HEADER = """\
import sys
if sys.platform == "win32":
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")
"""


class OfficeHandlerExpert:
    """Generates Office documents by having LLM produce Python scripts, then executing them."""

    def __init__(self, llm=None, tool_executor=None):
        self.llm = llm
        self.tools = tool_executor

    def run(self, ctx: TaskContext) -> dict:
        if not self.llm or not self.tools:
            return {
                "passed": False,
                "error": "OfficeHandler requires llm and tool_executor.",
            }

        # 1. Detect file type
        file_type = self._detect_type(ctx.user_input)

        # 2. Select scene system prompt
        system_map = {
            "xlsx": XLSX_SYSTEM,
            "pptx": PPTX_SYSTEM,
            "docx": DOCX_SYSTEM,
        }
        system = system_map[file_type]

        # Prepend expert_system_prompt from registry if available
        expert_prompt = ctx.expert_system_prompt or ""
        if expert_prompt:
            system = f"{expert_prompt}\n\n{system}"

        # 3. Determine output path
        output_path = self._get_output_path(ctx.user_input, file_type)

        # 4. Ask LLM to generate the script
        prompt = (
            f"用户需求：{ctx.user_input}\n\n"
            f"输出文件路径：{output_path}\n\n"
            f"生成完整的Python脚本，把文件保存到上面的路径。"
        )

        raw = self.llm.generate(
            prompt=prompt,
            system=system,
            max_tokens=4096,
            temperature=0.3,
        )

        script = self._extract_code(raw)
        if not script or len(script.strip()) < 30:
            return {
                "passed": False,
                "error": "LLM未生成有效的Python脚本",
            }

        # 5. Ensure script has UTF-8 header
        if "sys.stdout.reconfigure" not in script:
            script = _SCRIPT_HEADER + "\n" + script

        # 5.5 Auto-fix common LLM typos and syntax check
        script = self._auto_fix_script(script)
        syntax_err = self._syntax_check(script)
        if syntax_err:
            logger.warning("Office script syntax error after auto-fix: %s", syntax_err[:200])
            return {
                "passed": False,
                "error": f"生成的脚本有语法错误：{syntax_err[:300]}",
            }

        # 6. Write to temp file and execute
        tmp_fd, tmp_path = tempfile.mkstemp(suffix=".py", prefix="kwcode_office_")
        try:
            with os.fdopen(tmp_fd, "w", encoding="utf-8") as f:
                f.write(script)

            # Ensure output directory exists
            output_dir = os.path.dirname(output_path)
            if output_dir:
                os.makedirs(output_dir, exist_ok=True)

            stdout, stderr, rc = self.tools.run_bash(
                f'python "{tmp_path}"',
                cwd=ctx.project_root,
                timeout=60,
            )

            success = os.path.exists(output_path)

            if success:
                # Store result for orchestrator
                ctx.generator_output = {
                    "patches": [],
                    "explanation": f"已生成：{output_path}",
                }
                return {
                    "passed": True,
                    "output_path": output_path,
                    "output": f"已生成：{output_path}",
                }
            else:
                error_msg = stderr.strip() or stdout.strip() or "脚本执行完毕但文件未生成"
                logger.warning("Office script failed: rc=%d, stderr=%s", rc, error_msg[:200])
                return {
                    "passed": False,
                    "error": f"脚本执行失败：{error_msg[:300]}",
                }
        finally:
            try:
                os.unlink(tmp_path)
            except OSError:
                pass

    @staticmethod
    def _detect_type(user_input: str) -> str:
        """Detect target file type from user input."""
        text = user_input.lower()
        if any(k in text for k in ["excel", "xlsx", "表格", "报表", "电子表格", "数据表", "财务"]):
            return "xlsx"
        if any(k in text for k in ["ppt", "pptx", "演示", "幻灯", "slide", "deck", "汇报材料"]):
            return "pptx"
        return "docx"

    @staticmethod
    def _get_output_path(user_input: str, file_type: str) -> str:
        """Determine output file path. Defaults to Desktop."""
        desktop = os.path.join(os.path.expanduser("~"), "Desktop")
        if not os.path.isdir(desktop):
            desktop = os.path.expanduser("~")

        # Try to extract meaningful name from user input
        chinese = re.findall(r"[\u4e00-\u9fff]+", user_input)
        if chinese:
            # Take first meaningful Chinese phrase, cap at 8 chars
            name = "".join(chinese)[:8]
        else:
            name = f"output_{int(time.time())}"

        # Sanitize filename
        name = re.sub(r'[\\/:*?"<>|]', "_", name)

        full_path = os.path.join(desktop, f"{name}.{file_type}")

        # Avoid overwriting existing files
        if os.path.exists(full_path):
            base = name
            for i in range(1, 100):
                candidate = os.path.join(desktop, f"{base}_{i}.{file_type}")
                if not os.path.exists(candidate):
                    return candidate

        return full_path

    @staticmethod
    def _extract_code(text: str) -> str:
        """Extract Python code from LLM output, stripping markdown fences and thinking tags."""
        text = text.strip()
        # Strip <think>...</think> blocks
        text = re.sub(r"<think>.*?</think>", "", text, flags=re.DOTALL).strip()
        # Extract from ```python ... ``` blocks
        match = re.search(r"```(?:python)?\s*\n(.*?)```", text, re.DOTALL)
        if match:
            return match.group(1).strip()
        # If no code block, return as-is (LLM followed instructions)
        return text.strip()

    @staticmethod
    def _auto_fix_script(script: str) -> str:
        """Fix common LLM typos in generated Python scripts."""
        # Double equals in assignment: `= = True` -> `= True`
        script = re.sub(r"=\s*=\s*True", "= True", script)
        script = re.sub(r"=\s*=\s*False", "= False", script)
        script = re.sub(r"=\s*=\s*None", "= None", script)
        # Triple quotes accidentally broken
        script = script.replace("'' '", "'''").replace("' ''", "'''")
        return script

    @staticmethod
    def _syntax_check(script: str) -> str:
        """Compile-check the script. Returns error string or empty if OK."""
        try:
            compile(script, "<office_script>", "exec")
            return ""
        except SyntaxError as e:
            return f"Line {e.lineno}: {e.msg}"
